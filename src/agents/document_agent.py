from .base_agent import BaseAgent
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
from typing import Dict, Any
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from src.connectors.standard_input import StandardInput
from src.preprocessing.engine import PreprocessingEngine

class DocumentAgent(BaseAgent):
    def __init__(self, event_bus, working_memory):
        super().__init__(
            agent_id="document_agent",
            capabilities=["parse_pdf", "parse_docx", "parse_xlsx", "parse_txt", "parse_pptx"],
            event_bus=event_bus,
            working_memory=working_memory
        )
        self.preprocessor = PreprocessingEngine()
    
    def process(self, task: Dict[str, Any]) -> Dict[str, Any]:
        file_path = task.get("file_path")
        file_type = task.get("file_type")
        
        print(f"[DocumentAgent] Processing {file_type} file: {file_path}")
        
        if file_type == "pdf":
            text = self._parse_pdf(file_path)
        elif file_type == "docx":
            text = self._parse_docx(file_path)
        elif file_type == "xlsx":
            text = self._parse_xlsx(file_path)
        elif file_type == "txt":
            text = self._parse_txt(file_path)
        elif file_type == "pptx":
            text = self._parse_pptx(file_path)
        else:
            return {"error": "Unsupported file type"}
        
        print(f"[DocumentAgent] Extracted {len(text)} characters")
        
        # Create standardized input
        std_input = StandardInput.create(
            source_type="file",
            raw_text=text,
            metadata={"file_path": file_path, "file_type": file_type}
        )
        
        # Preprocess
        processed = self.preprocessor.process(std_input.raw_text, std_input.metadata)
        
        result = {
            "input_id": std_input.id,
            "file_path": file_path,
            "text": processed["cleaned_text"],
            "deadlines": processed["deadlines"],
            "chunks": processed["chunks"],
            "confidence": 1.0
        }
        
        print(f"[DocumentAgent] Publishing document_processed event")
        self.publish_event("document_processed", result)
        return result
    
    def _parse_pdf(self, file_path: str) -> str:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text
    
    def _parse_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    
    def _parse_xlsx(self, file_path: str) -> str:
        wb = openpyxl.load_workbook(file_path)
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell]) + "\n"
        return text
    
    def _parse_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _parse_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
