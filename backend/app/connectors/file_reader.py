"""
File reader for various document formats.
Supports PDF, DOCX, XLSX, PPTX, TXT.
"""
from pathlib import Path
from typing import Optional


def read_pdf(file_path: Path) -> str:
    """Read PDF file and extract text."""
    try:
        from PyPDF2 import PdfReader  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "PyPDF2 is not installed. Install optional AI/doc dependencies with:\n"
            "  pip install -r backend/requirements-ai.txt"
        ) from exc

    try:
        reader = PdfReader(str(file_path))
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text() or "")
        return "\n".join(texts)
    except Exception as e:  # pragma: no cover - I/O heavy
        raise Exception(f"Error reading PDF file {file_path}: {e}")


def read_docx(file_path: Path) -> str:
    """Read DOCX file and extract text."""
    try:
        import docx  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "python-docx is not installed. Install optional AI/doc dependencies with:\n"
            "  pip install -r backend/requirements-ai.txt"
        ) from exc

    try:
        doc = docx.Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:  # pragma: no cover - I/O heavy
        raise Exception(f"Error reading DOCX file {file_path}: {e}")


def read_xlsx(file_path: Path) -> str:
    """Read XLSX file and extract text."""
    try:
        import pandas as pd  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "pandas/openpyxl are not installed. Install optional AI/doc dependencies with:\n"
            "  pip install -r backend/requirements-ai.txt"
        ) from exc

    try:
        # Read all sheets and concatenate as text
        xls = pd.ExcelFile(str(file_path))
        parts = []
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            parts.append(f"=== Sheet: {sheet_name} ===")
            parts.append(df.to_string(index=False))
        return "\n\n".join(parts)
    except Exception as e:  # pragma: no cover - I/O heavy
        raise Exception(f"Error reading XLSX file {file_path}: {e}")


def read_pptx(file_path: Path) -> str:
    """Read PPTX file and extract text."""
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "python-pptx is not installed. Install optional AI/doc dependencies with:\n"
            "  pip install python-pptx"
        ) from exc

    try:
        prs = Presentation(str(file_path))
        texts = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts.append(f"=== Slide {slide_num} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n\n".join(texts)
    except Exception as e:  # pragma: no cover - I/O heavy
        raise Exception(f"Error reading PPTX file {file_path}: {e}")


def read_txt(file_path: Path) -> str:
    """Read TXT file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:  # pragma: no cover - I/O heavy
        raise Exception(f"Error reading TXT file {file_path}: {e}")


def read_file(file_path: Path) -> Optional[str]:
    """Read file based on extension."""
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return read_pdf(file_path)
    if ext == ".docx":
        return read_docx(file_path)
    if ext == ".xlsx":
        return read_xlsx(file_path)
    if ext == ".pptx":
        return read_pptx(file_path)
    if ext == ".txt":
        return read_txt(file_path)

    raise ValueError(f"Unsupported file type: {ext}")

